from __future__ import annotations

import csv
from io import StringIO
from pathlib import Path
from typing import List

from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain_core.documents import Document

try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:  # pragma: no cover - older / partial venvs
    try:
        from langchain.text_splitter import RecursiveCharacterTextSplitter
    except ImportError as e:
        raise ImportError(
            "Install text splitters: pip install langchain-text-splitters "
            "(or run: pip install -r requirements.txt)"
        ) from e

from src import config

# Plain-text-like (CSV/TSV handled separately for table parsing)
_TEXT_LIKE = {".txt", ".md", ".markdown", ".log", ".rst", ".json", ".xml", ".html", ".htm"}


def _doc_from_text(text: str, path: Path, *, file_name: str | None = None) -> Document:
    body = (text or "").strip()
    if not body:
        body = "(empty document)"
    label = file_name or path.name
    return Document(page_content=body, metadata={"source": str(path), "file_name": label})


def _text_from_docx(path: Path) -> str:
    try:
        from docx import Document as DocxDocument
    except ModuleNotFoundError as e:
        raise ImportError(
            "Word .docx files need the **python-docx** package. From your venv run: "
            "`pip install python-docx` or `pip install -r requirements.txt`."
        ) from e

    doc = DocxDocument(str(path))
    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        rows_txt: list[str] = []
        for row in table.rows:
            rows_txt.append("\t".join(cell.text.replace("\n", " ") for cell in row.cells))
        if rows_txt:
            parts.append("\n".join(rows_txt))
    return "\n\n".join(parts)


def _text_from_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ModuleNotFoundError as e:
        raise ImportError(
            "Excel .xlsx needs **openpyxl**. Run: `pip install openpyxl` or "
            "`pip install -r requirements.txt`."
        ) from e

    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    lines: list[str] = []
    try:
        for sheet in wb.worksheets:
            lines.append(f"## Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                line = "\t".join("" if v is None else str(v) for v in row)
                if line.strip():
                    lines.append(line)
    finally:
        wb.close()
    return "\n".join(lines)


def _text_from_xls(path: Path) -> str:
    try:
        import xlrd
    except ModuleNotFoundError as e:
        raise ImportError(
            "Legacy Excel .xls needs **xlrd**. Run: `pip install xlrd` or "
            "`pip install -r requirements.txt`."
        ) from e

    book = xlrd.open_workbook(str(path))
    lines: list[str] = []
    for si in range(book.nsheets):
        sh = book.sheet_by_index(si)
        lines.append(f"## Sheet: {sh.name}")
        for ri in range(sh.nrows):
            vals = sh.row_values(ri)
            line = "\t".join("" if v is None else str(v) for v in vals)
            if line.strip():
                lines.append(line)
    return "\n".join(lines)


def _text_from_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as e:
        raise ImportError(
            "PowerPoint .pptx needs **python-pptx**. Run: `pip install python-pptx` or "
            "`pip install -r requirements.txt`."
        ) from e

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide_parts:
            parts.append(f"## Slide {i}\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)


def _text_from_csv(path: Path, delimiter: str | None = None) -> str:
    text: str | None = None
    for enc in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            text = path.read_text(encoding=enc)
            break
        except UnicodeDecodeError:
            continue
    if text is None:
        text = path.read_text(encoding="utf-8", errors="replace")

    sample = text[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=delimiter or ",\t;")
        delim = dialect.delimiter
    except csv.Error:
        delim = delimiter or ","

    reader = csv.reader(StringIO(text), delimiter=delim)
    lines = ["\t".join(row) for row in reader]
    return "\n".join(lines)


class DocumentLoader:
    def __init__(self, chunk_size: int | None = None, chunk_overlap: int | None = None):
        self.chunk_size = chunk_size if chunk_size is not None else config.DOC_CHUNK_SIZE
        self.chunk_overlap = chunk_overlap if chunk_overlap is not None else config.DOC_CHUNK_OVERLAP
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
        )

    def load_documents(self, file_path: str, *, display_name: str | None = None) -> List[Document]:
        path = Path(file_path)
        suffix = path.suffix.lower()
        label = (display_name or path.name).strip() or path.name

        if suffix == ".pdf":
            loader = PyPDFLoader(str(path))
            documents = loader.load()
        elif suffix == ".doc":
            raise ValueError(
                "Legacy Word .doc is not supported. Save the file as **.docx** or **PDF**, then upload again."
            )
        elif suffix == ".docx":
            documents = [_doc_from_text(_text_from_docx(path), path, file_name=label)]
        elif suffix == ".xlsx":
            documents = [_doc_from_text(_text_from_xlsx(path), path, file_name=label)]
        elif suffix == ".xls":
            documents = [_doc_from_text(_text_from_xls(path), path, file_name=label)]
        elif suffix == ".pptx":
            documents = [_doc_from_text(_text_from_pptx(path), path, file_name=label)]
        elif suffix == ".tsv":
            documents = [_doc_from_text(_text_from_csv(path, delimiter="\t"), path, file_name=label)]
        elif suffix == ".csv":
            documents = [_doc_from_text(_text_from_csv(path), path, file_name=label)]
        elif suffix in _TEXT_LIKE:
            enc = "utf-8-sig" if suffix in (".csv", ".tsv", ".json", ".xml") else "utf-8"
            loader = TextLoader(str(path), encoding=enc)
            documents = loader.load()
        else:
            raise ValueError(
                f"Unsupported file type: {suffix}. "
                f"Supported: PDF, TXT, Markdown, CSV/TSV, DOCX, XLSX, XLS, PPTX."
            )

        for d in documents:
            meta = dict(d.metadata or {})
            meta["source"] = str(path)
            meta["file_name"] = label
            d.metadata = meta

        return self.text_splitter.split_documents(documents)

    def load_from_folder(self, folder_path: str) -> List[Document]:
        all_docs: List[Document] = []
        folder = Path(folder_path)
        seen: set[Path] = set()

        patterns = (
            "**/*.pdf",
            "**/*.txt",
            "**/*.md",
            "**/*.markdown",
            "**/*.csv",
            "**/*.tsv",
            "**/*.docx",
            "**/*.xlsx",
            "**/*.xls",
            "**/*.pptx",
            "**/*.json",
            "**/*.xml",
            "**/*.html",
            "**/*.htm",
        )
        for pattern in patterns:
            for file_path in folder.glob(pattern):
                if file_path in seen:
                    continue
                seen.add(file_path)
                try:
                    all_docs.extend(self.load_documents(str(file_path)))
                except Exception as e:
                    print(f"Error loading {file_path}: {e}")

        return all_docs
